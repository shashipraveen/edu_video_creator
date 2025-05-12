from pptx import Presentation
from pptx.util import Inches

def create_ppt(content, output_file):
    """
    Creates a PowerPoint presentation from the provided content.
    """
    prs = Presentation()

    for slide_title, slide_content in content.items():
        slide_layout = prs.slide_layouts[1]  # Using a layout with title and content
        slide = prs.slides.add_slide(slide_layout)

        title = slide.shapes.title
        title.text = slide_title

        content_box = slide.shapes.placeholders[1]
        content_box.text = slide_content

    prs.save(output_file)
    print(f"PowerPoint saved to {output_file}")

# Example usage
if __name__ == "__main__":
    content = {
        "Introduction to Atoms": "Atoms are the basic building blocks of matter.",
        "Atomic Structure": "Each atom consists of a nucleus containing protons and neutrons, surrounded by electrons."
    }
    create_ppt(content, "output_presentation.pptx")
